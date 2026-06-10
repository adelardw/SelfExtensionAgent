from pathlib import Path
from langchain_core.tools import tool

@tool
def create_evolution_pptx() -> str:
    """Создаёт презентацию .pptx по философии эволюции кибернетики и сохраняет её на диск.
    
    Returns:
        str: Путь к созданному файлу и подтверждение.
    """
    # Import pptx inside the function
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return "Ошибка: библиотека python-pptx не установлена. Установите её командой: pip install python-pptx"
    
    output_path = "/Users/Shared/evolution_cybernetics_philosophy.pptx"
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
    DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    def add_line(slide, left, top, width):
        shape = slide.shapes.add_shape(1, left, top, width, Pt(2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_BLUE
        shape.line.fill.background()

    def add_title_bar(slide, title_text):
        txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.font.name = "Calibri"
        add_line(slide, Inches(1), Inches(1.3), Inches(11.333))

    def add_body_text(slide, text, top):
        txBox = slide.shapes.add_textbox(Inches(1.2), top, Inches(11), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        lines = text.strip().split('\n')
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            line = line.strip()
            if line:
                p.text = line
                p.font.size = Pt(16)
                p.font.color.rgb = DARK_GRAY
                p.font.name = "Calibri"
                p.space_after = Pt(8)

    # S1 - Title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    b = s.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.333), Inches(1.5))
    tf = b.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Эволюция кибернетики:"; p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = DARK_BLUE; p.font.name = "Calibri"; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = "философский взгляд"; p2.font.size = Pt(40); p2.font.bold = True; p2.font.color.rgb = DARK_BLUE; p2.font.name = "Calibri"; p2.alignment = PP_ALIGN.CENTER
    add_line(s, Inches(4), Inches(3.7), Inches(5.333))
    b2 = s.shapes.add_textbox(Inches(1.5), Inches(4.0), Inches(10.333), Inches(0.8))
    tf2 = b2.text_frame; tf2.word_wrap = True
    p3 = tf2.paragraphs[0]; p3.text = "От античных автоматов к самоорганизующимся системам"; p3.font.size = Pt(22); p3.font.color.rgb = DARK_GRAY; p3.font.name = "Calibri"; p3.alignment = PP_ALIGN.CENTER
    b3 = s.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(10.333), Inches(0.5))
    tf3 = b3.text_frame; tf3.word_wrap = True
    p4 = tf3.paragraphs[0]; p4.text = "Философский семинар, 2024"; p4.font.size = Pt(16); p4.font.color.rgb = RGBColor(0x99,0x99,0x99); p4.font.name = "Calibri"; p4.alignment = PP_ALIGN.CENTER

    # S2
    s = prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    add_title_bar(s, "Что такое кибернетика?")
    add_body_text(s, "• Норберт Винер (1948): «Кибернетика — наука об управлении и связи в машинах и живых организмах»\n• От греч. κυβερνήτης (kybernētēs) — «кормчий», «управляющий кораблём»\n• Ключевые понятия: обратная связь, гомеостазис, саморегуляция, информация\n• Философский смысл: кибернетика стирает границу между живым и неживым, предлагая единый язык описания", Inches(1.8))

    # S3
    s = prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    add_title_bar(s, "Предыстория: от античности до Просвещения")
    add_body_text(s, "• Аристотель: идея целевой причины (teleological causation) — предвосхищение обратной связи\n• Герон Александрийский (I в. н.э.): автоматы и механизмы с простейшей регуляцией\n• Рене Декарт: животные как «сложные машины» (mechanical philosophy)\n• Жюльен Ламетри (1747): «Человек-машина» — материалистическая философия\n• Адам Смит: «невидимая рука рынка» как пример саморегулирующейся системы", Inches(1.8))

    # S4
    s = prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    add_title_bar(s, "Первая кибернетика (1940–1960)")
    add_body_text(s, "• Норберт Винер: «Кибернетика» (1948) — рождение науки\n• У. Росс Эшби: «Введение в кибернетику» (1956) — закон необходимого разнообразия\n• Джон фон Нейман: теория клеточных автоматов, самовоспроизводящиеся машины\n• Философский смысл: телеология без теологии — целенаправленное поведение как результат обратной связи\n• Ключевой вопрос: может ли машина имитировать жизнь?", Inches(1.8))

    # S5
    s = prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    add_title_bar(s, "Вторая кибернетика: наблюдатель включён в систему")
    add_body_text(s, "• Хайнц фон Фёрстер: кибернетика кибернетики (second-order cybernetics)\n• Ключевой сдвиг: наблюдатель больше не отделён от системы — он её часть\n• Умберто Матурана и Франсиско Варела: автопоэзис (аутопоэзис) — самопорождение живых систем\n• Философский смысл: пересмотр объективности — познание не отражает реальность, а конструирует её\n• Радикальный конструктивизм: мир, который мы познаём, есть мир, который мы создаём", Inches(1.8))

    # S6
    s = prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    add_title_bar(s, "Кибернетика третьего порядка")
    add_body_text(s, "• Расширение за пределы биологии: социальные системы, экономика, культура\n• Никлас Луман: социальные системы как аутопоэтические системы коммуникации\n• Йозеф Вандри: семиотическая кибернетика — знаки и значения как управляющие сигналы\n• Философский смысл: общество функционирует как кибернетическая система с обратными связями\n• Двойная контингентность (Луман): каждое действие зависит от ожидания реакции другого", Inches(1.8))

    # S7
    s = prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    add_title_bar(s, "Кибернетика и искусственный интеллект")
    add_body_text(s, "• Связь кибернетики с ИИ: от перцептронов (Розенблатт, 1958) до нейросетей\n• Отличия: кибернетика → управление и целостность; ИИ → вычисления и алгоритмы\n• Символический ИИ vs коннекционизм: старый спор о природе разума\n• Философский вопрос: может ли машина обладать сознанием?\n• Тезис «Китайской комнаты» (Сёрль, 1980): синтаксис ≠ семантика\n• Современный синтез: глубокое обучение как кибернетическая система с обратной связью", Inches(1.8))

    # S8
    s = prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    add_title_bar(s, "Философские вызовы и этика")
    add_body_text(s, "• Проблема свободы воли: если всё — кибернетические системы, существует ли свобода?\n• Ответственность: кто отвечает за действия автономной системы?\n• Прозрачность: «чёрный ящик» нейросетей — невозможность объяснения решений\n• Биоэтика: киборгизация человека, нейроинтерфейсы, трансгуманизм\n• Социальная кибернетика: риски тотального контроля (антиутопии XX века)\n• Критический вопрос: кибернетика — инструмент освобождения или порабощения?", Inches(1.8))

    # S9
    s = prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    add_title_bar(s, "Выводы")
    add_body_text(s, "• Кибернетика прошла путь от механицизма к конструктивизму и обратно к синтезу\n• Каждый этап эволюции кибернетики ставил фундаментальные философские вопросы\n• Три волны: управление → наблюдение → самоорганизация\n• Кибернетика не просто описывает мир — она предлагает новый способ мышления\n• Современное значение: кибернетическое мышление необходимо для понимания ИИ, биотехнологий и глобальных систем\n• Финальный тезис: эволюция кибернетики — это эволюция самого понятия «знание»", Inches(1.8))

    # S10
    s = prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    add_title_bar(s, "Источники и литература")
    add_body_text(s, "• Винер Н. «Кибернетика, или Управление и связь в животном и машине» (1948)\n• Эшби У. Р. «Введение в кибернетику» (1956)\n• Фёрстер Х. фон. «Кибернетика кибернетики» (1974)\n• Матурана У., Варела Ф. «Древо познания» (1984)\n• Луман Н. «Социальные системы» (1984)\n• Сёрль Дж. «Разум, мозг и наука» (1984)\n• Винер Н. «Кибернетика и общество» (1950)", Inches(1.8))

    # Save
    prs.save(output_path)
    file_size = Path(output_path).stat().st_size
    
    return f"SUCCESS: Презентация создана!\nФайл: {output_path}\nСлайдов: {len(prs.slides)}\nРазмер: {file_size} байт"