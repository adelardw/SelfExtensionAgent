"""
Работа с медиа-вложениями: картинки (vision) и голос (транскрипция).

Без отдельных платных STT/vision-сервисов: gemini-2.5-flash-lite мультимодален
(изображения + аудио на входе) и уже является fast-тиром агента — то есть
вложения стоят те же копейки, что и обычный запрос.

  describe_image(path, question)  → текстовое описание/анализ изображения
  transcribe_audio(path)          → расшифровка голосового сообщения
  attachment_context(paths, q)    → готовый блок контекста по списку вложений
"""
from __future__ import annotations

import base64
import mimetypes
import shutil
import subprocess
import tempfile
from pathlib import Path

from langchain_core.messages import HumanMessage

from .llm import chat

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".heic"}
AUDIO_EXTS = {".ogg", ".oga", ".mp3", ".wav", ".m4a", ".aiff", ".flac"}
VIDEO_EXTS = {".mp4", ".mov", ".avi", ".mkv", ".webm", ".gif"}  # кадры→vision + аудио→транскрипт
TEXT_EXTS = {".txt", ".md", ".csv", ".json", ".yml", ".yaml", ".py", ".log", ".html", ".xml"}
DOC_EXTS = {".pdf", ".xlsx", ".xls", ".docx", ".pptx"}  # документы — извлекаем текст/таблицы

MAX_INLINE_TEXT = 6000  # сколько символов текстового файла инлайнить в контекст


def _pdf_pymupdf(path: Path, max_chars: int) -> str:
    import fitz  # pymupdf

    doc = fitz.open(str(path))
    parts = []
    for page in doc:
        parts.append(page.get_text())
        if sum(len(p) for p in parts) > max_chars:
            break
    doc.close()
    return "\n".join(parts)[:max_chars]


def _pdf_liteparse(path: Path, max_chars: int) -> str:
    """liteparse: сохраняет ЛЕЙАУТ (строки таблиц выровнены) — LLM читает «значение↔строка»."""
    import liteparse

    # OCR off — для текстовых PDF быстро (~0.04с); включаем только если текста нет (скан).
    text = liteparse.LiteParse(ocr_enabled=False, quiet=True).parse(str(path)).text
    if len((text or "").strip()) < 20:  # вероятно скан → пробуем OCR
        text = liteparse.LiteParse(ocr_enabled=True, quiet=True).parse(str(path)).text
    return (text or "")[:max_chars]


def _pdf_opendataloader(path: Path, max_chars: int) -> str:
    """opendataloader (#1 точность таблиц/LaTeX) — ТРЕБУЕТ Java 11+; опциональный power-тир."""
    import opendataloader_pdf

    out = opendataloader_pdf.convert(input_path=[str(path)], format="markdown")
    return (out if isinstance(out, str) else str(out))[:max_chars]


def read_pdf(path: Path, max_chars: int = 12000) -> str:
    """
    PDF → текст тиерами (по убыванию точности таблиц, с фолбэком на надёжность):
      1. liteparse — layout-aware (строки таблиц выровнены), быстро, без JVM — ДЕФОЛТ;
      2. opendataloader — #1 точность (таблицы/LaTeX), но требует Java — только если включён
         AGENT_PDF_POWER=1 и Java есть;
      3. pymupdf — простой надёжный baseline.
    """
    import os

    if os.getenv("AGENT_PDF_POWER") == "1":
        try:
            t = _pdf_opendataloader(path, max_chars)
            if t.strip():
                return t
        except Exception:  # noqa: BLE001 — нет Java/ошибка → следующий тир
            pass
    try:
        t = _pdf_liteparse(path, max_chars)
        if t.strip():
            return t
    except Exception:  # noqa: BLE001
        pass
    return _pdf_pymupdf(path, max_chars)


def read_excel(path: Path, max_chars: int = 12000) -> str:
    """Все листы Excel как текстовые таблицы (pandas)."""
    import pandas as pd

    out = []
    for name, df in pd.read_excel(str(path), sheet_name=None).items():
        out.append(f"[Лист: {name}]\n{df.to_string(max_rows=200)}")
    return "\n\n".join(out)[:max_chars]


def read_docx(path: Path, max_chars: int = 12000) -> str:
    """Текст из .docx (python-docx)."""
    import docx

    d = docx.Document(str(path))
    return "\n".join(p.text for p in d.paragraphs)[:max_chars]


def read_pptx(path: Path, max_chars: int = 12000) -> str:
    """Текст всех слайдов .pptx (python-pptx): заголовки, буллеты, текст фигур, заметки."""
    from pptx import Presentation

    prs = Presentation(str(path))
    out = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [shape.text.strip() for shape in slide.shapes
                 if shape.has_text_frame and shape.text.strip()]
        if parts:
            out.append(f"[Слайд {i}]\n" + "\n".join(parts))
    return "\n\n".join(out)[:max_chars] or "[презентация без текста]"


def _ffmpeg_frames(path: Path, n: int = 3) -> list[Path]:
    """Сэмпл до n кадров из видео/gif через ffmpeg (равномерно по времени)."""
    if not shutil.which("ffmpeg"):
        return []
    out_dir = Path(tempfile.mkdtemp())
    # один кадр раз в ~5с, максимум n — дёшево и репрезентативно для «о чём ролик»
    pat = str(out_dir / "f_%02d.png")
    try:
        subprocess.run(["ffmpeg", "-v", "error", "-i", str(path), "-vf", "fps=1/5",
                        "-frames:v", str(n), pat], check=True, timeout=60)
    except Exception:  # noqa: BLE001
        return []
    return sorted(out_dir.glob("f_*.png"))


def read_video(path: Path, question: str = "") -> str:
    """
    Универсальный разбор видео/gif БЕЗ спец-сервисов: сэмпл кадров → vision-описание +
    транскрипт аудио-дорожки (если есть). Кадры и звук стоят те же копейки (fast-тир).
    """
    frames = _ffmpeg_frames(path, n=3)
    parts = []
    for fr in frames:
        try:
            parts.append(describe_image(fr, question))
        except Exception:  # noqa: BLE001
            pass
    visual = "\n".join(f"- кадр {i+1}: {d}" for i, d in enumerate(parts)) if parts \
        else "(кадры не извлечены — нужен ffmpeg)"
    # аудио-дорожка → транскрипт (видео часто без звука — тогда пропускаем)
    audio_txt = ""
    if shutil.which("ffmpeg") and path.suffix.lower() != ".gif":
        wav = Path(tempfile.mktemp(suffix=".mp3"))
        try:
            subprocess.run(["ffmpeg", "-v", "error", "-i", str(path), "-vn", "-ac", "1", str(wav)],
                           check=True, timeout=120)
            if wav.exists() and wav.stat().st_size > 1000:
                audio_txt = transcribe_audio(wav)
        except Exception:  # noqa: BLE001
            pass
    res = f"[Видео {path.name}] Кадры:\n{visual}"
    if audio_txt:
        res += f"\n\n[Аудио-дорожка]:\n{audio_txt}"
    return res


def read_file(path: str | Path, max_chars: int = 12000) -> str:
    """
    Универсальный ридер вложения по расширению: документ→текст/таблицы, картинка→vision,
    аудио→транскрипт, текст→содержимое. Используется attachment_context и навыком.
    """
    p = Path(path)
    if not p.exists():
        return f"[файл не найден: {p}]"
    ext = p.suffix.lower()
    try:
        if ext == ".pdf":
            return read_pdf(p, max_chars)
        if ext in (".xlsx", ".xls"):
            return read_excel(p, max_chars)
        if ext == ".docx":
            return read_docx(p, max_chars)
        if ext == ".pptx":
            return read_pptx(p, max_chars)
        if ext in VIDEO_EXTS:        # видео/gif — раньше image (gif шёл одним кадром)
            return read_video(p)
        if ext in IMAGE_EXTS:
            return describe_image(p)
        if ext in AUDIO_EXTS:
            return transcribe_audio(p)
        return p.read_text(encoding="utf-8", errors="ignore")[:max_chars]
    except Exception as e:  # noqa: BLE001
        return f"[не смог прочитать {p.name}: {type(e).__name__}: {e}]"


def _b64(path: Path) -> str:
    return base64.b64encode(path.read_bytes()).decode()


def describe_image(path: str | Path, question: str = "") -> str:
    """Vision-анализ изображения fast-моделью. question — на что обратить внимание."""
    p = Path(path)
    mime = mimetypes.guess_type(p.name)[0] or "image/png"
    prompt = (
        "Опиши это изображение подробно и структурно: что на нём, весь видимый текст "
        "(перепиши точно), числа, таблицы, UI-элементы. Это описание будет контекстом "
        "для агента, отвечающего на запрос пользователя."
    )
    if question:
        prompt += f"\nЗапрос пользователя (учти при описании): {question}"
    msg = HumanMessage(content=[
        {"type": "text", "text": prompt},
        {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{_b64(p)}"}},
    ])
    resp = chat("fast").invoke([msg])
    return resp.content if hasattr(resp, "content") else str(resp)


def read_pdf_visual(path: str | Path, question: str = "", max_pages: int = 5) -> str:
    """
    VISION-чтение PDF: рендер страниц в картинки → описание ФИГУР/графиков/диаграмм через
    vision. Закрывает пробел, где ответ в ФИГУРЕ (оси, метки, графики), а не в тексте
    (GAIA L2/L3 с фигурами). Текстовый _pdf_pymupdf фигуры не видит — этот читает визуально.
    """
    import fitz  # pymupdf

    doc = fitz.open(str(path))
    parts = []
    for i in range(min(max_pages, doc.page_count)):
        page = doc.load_page(i)
        pix = page.get_pixmap(dpi=140)  # достаточно для чтения меток на осях/графиках
        tmp = Path(tempfile.mkstemp(suffix=".png")[1])
        pix.save(str(tmp))
        try:
            focus = (question or "") + " Особое внимание: ФИГУРЫ, графики, диаграммы — оси, "\
                    "метки на концах осей, подписи, числа, легенды. Перепиши их точно."
            parts.append(f"[Стр. {i + 1}]\n" + describe_image(tmp, focus))
        except Exception as e:  # noqa: BLE001
            parts.append(f"[Стр. {i + 1}: vision не сработал: {e}]")
        finally:
            try:
                tmp.unlink()
            except Exception:  # noqa: BLE001
                pass
    doc.close()
    return "\n\n".join(parts) or "[PDF без страниц]"


def make_pdf_vision_tool():
    """Тул read_pdf_figures: VISION-чтение фигур/графиков в PDF (ответ в картинке, не тексте)."""
    from langchain_core.tools import StructuredTool
    from pydantic import BaseModel, Field

    class _Q(BaseModel):
        path: str = Field(description="Путь к PDF-файлу")
        question: str = Field(description="Что искать в фигурах/графиках (метки осей, числа…)", default="")

    def _run(path: str, question: str = "") -> str:
        try:
            return read_pdf_visual(path, question)
        except Exception as e:  # noqa: BLE001
            return f"[не удалось прочитать фигуры PDF: {type(e).__name__}: {e}]"

    return StructuredTool.from_function(
        func=_run, name="read_pdf_figures", args_schema=_Q,
        description="VISION-read FIGURES/charts/diagrams in a PDF (axis labels, values, legends) when "
                    "the answer is in an image, not the text. Use after plain text extraction misses it.",
    )


def _to_mp3(path: Path) -> Path:
    """Конвертация в mp3 через ffmpeg (telegram-голосовые приходят в ogg/opus)."""
    if not shutil.which("ffmpeg"):
        return path  # попробуем отдать как есть
    out = Path(tempfile.mkstemp(suffix=".mp3")[1])
    res = subprocess.run(
        ["ffmpeg", "-y", "-i", str(path), "-ac", "1", "-b:a", "48k", str(out)],
        capture_output=True, timeout=60,
    )
    return out if res.returncode == 0 and out.stat().st_size > 0 else path


def transcribe_audio(path: str | Path) -> str:
    """Расшифровка аудио fast-моделью (audio-вход gemini). Возвращает чистый текст."""
    p = Path(path)
    fmt = p.suffix.lstrip(".").lower()
    if fmt not in ("mp3", "wav"):
        p = _to_mp3(p)
        fmt = p.suffix.lstrip(".").lower() or "mp3"
    msg = HumanMessage(content=[
        {"type": "text", "text": "Расшифруй это голосовое сообщение. Верни ТОЛЬКО текст сказанного, без комментариев."},
        {"type": "input_audio", "input_audio": {"data": _b64(p), "format": fmt}},
    ])
    resp = chat("fast").invoke([msg])
    return (resp.content if hasattr(resp, "content") else str(resp)).strip()


def attachment_context(paths: list[str | Path], question: str = "") -> str:
    """
    Готовит блок контекста по вложениям: картинки — vision-описание, текстовые файлы —
    инлайн содержимого (с обрезкой), прочее — путь+размер (deliberate-путь откроет
    навыками file_operations/text_file_processor).
    """
    blocks = []
    for raw in paths:
        p = Path(raw)
        if not p.exists():
            blocks.append(f"[Вложение {p.name}: файл не найден]")
            continue
        ext = p.suffix.lower()
        if ext in IMAGE_EXTS:
            try:
                desc = describe_image(p, question)
                blocks.append(f"[Изображение {p.name} (путь: {p})]\nОписание:\n{desc}")
            except Exception as e:  # noqa: BLE001
                blocks.append(f"[Изображение {p.name} (путь: {p}) — vision не сработал: {e}]")
        elif ext in DOC_EXTS:
            content = read_file(p, MAX_INLINE_TEXT)
            kind = {"pdf": "PDF", "xlsx": "Excel", "xls": "Excel", "docx": "Word",
                    "pptx": "Презентация"}.get(ext.lstrip("."), "Документ")
            blocks.append(f"[{kind} {p.name} (путь: {p})]\n{content}")
        elif ext in VIDEO_EXTS:
            try:
                blocks.append(f"[Видео {p.name} (путь: {p})]\n{read_video(p, question)}")
            except Exception as e:  # noqa: BLE001
                blocks.append(f"[Видео {p.name} — разбор не сработал: {e}]")
        elif ext in AUDIO_EXTS:
            try:
                blocks.append(f"[Аудио {p.name} (путь: {p})]\nТранскрипт:\n{transcribe_audio(p)}")
            except Exception as e:  # noqa: BLE001
                blocks.append(f"[Аудио {p.name} — транскрипт не сработал: {e}]")
        elif ext in TEXT_EXTS:
            try:
                text = p.read_text(encoding="utf-8", errors="ignore")
                cut = text[:MAX_INLINE_TEXT]
                tail = f"\n…(обрезано, всего {len(text)} симв.; полный файл: {p})" if len(text) > MAX_INLINE_TEXT else ""
                blocks.append(f"[Файл {p.name} (путь: {p})]\n{cut}{tail}")
            except Exception as e:  # noqa: BLE001
                blocks.append(f"[Файл {p.name} (путь: {p}) — не читается: {e}]")
        else:
            blocks.append(
                f"[Файл {p.name} (путь: {p}, {p.stat().st_size} байт) — бинарный/другой формат; "
                f"открой инструментами при необходимости]"
            )
    return "\n\n".join(blocks)
