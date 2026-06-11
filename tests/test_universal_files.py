"""Универсальность файлов: pptx-чтение + маршрутизация видео/gif. Офлайн."""
import tempfile
from pathlib import Path

import pytest


def test_read_pptx_extracts_text():
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2)).text_frame
    tb.text = "Эволюция киберагентов: тезис о харнессе"
    p = Path(tempfile.mktemp(suffix=".pptx"))
    prs.save(str(p))

    from src.media import read_file
    out = read_file(p)
    assert "харнессе" in out and "Слайд 1" in out


def test_video_and_gif_routed_to_video_not_image():
    from src.media import VIDEO_EXTS, IMAGE_EXTS
    assert ".gif" in VIDEO_EXTS and ".gif" not in IMAGE_EXTS   # gif → кадры, не один снимок
    assert ".mp4" in VIDEO_EXTS and ".mov" in VIDEO_EXTS


def test_pptx_in_doc_exts():
    from src.media import DOC_EXTS
    assert ".pptx" in DOC_EXTS
